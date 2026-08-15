"""PPTX slide-order adapter for text, tables, images, and speaker notes."""

from __future__ import annotations

import hashlib
import io
import time
from dataclasses import dataclass

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from semikb_ingest.assets import ProcessPayloadStore, inspect_image
from semikb_ingest.errors import IngestError, IngestErrorCode
from semikb_ingest.models import (
    CaptionSource,
    ImageAssetDraft,
    SourceFormat,
    SourceLocation,
    TableAssetDraft,
)
from semikb_ingest.parsers.common import (
    complete_document,
    matrix_to_html,
    matrix_to_markdown,
    normalized_text,
)
from semikb_ingest.parsers.registry import ParseRequest
from semikb_ingest.structure import BlockKind, StructuredBlock

_EMU_PER_POINT = 12_700


@dataclass(frozen=True, slots=True)
class PptxLimits:
    max_slides: int = 2_000
    max_shapes: int = 100_000
    max_images: int = 10_000


class PptxStructuredParser:
    parser_id = "pptx-structured-v1"
    parser_version = "1.0.0"
    source_format = SourceFormat.PPTX

    def __init__(
        self,
        payload_store: ProcessPayloadStore,
        limits: PptxLimits | None = None,
    ) -> None:
        self.payload_store = payload_store
        self.limits = limits or PptxLimits()

    def parse(self, request: ParseRequest):
        started = time.monotonic()
        try:
            presentation = Presentation(io.BytesIO(request.content))
        except Exception as exc:
            raise IngestError(
                IngestErrorCode.CORRUPT_DOCUMENT,
                "The PPTX presentation could not be opened.",
            ) from exc
        if len(presentation.slides) > self.limits.max_slides:
            self._raise_limit()

        blocks: list[StructuredBlock] = []
        images: list[ImageAssetDraft] = []
        tables: list[TableAssetDraft] = []
        image_by_hash: dict[str, ImageAssetDraft] = {}
        shape_count = 0
        detected_title: str | None = None

        for slide_number, slide in enumerate(presentation.slides, start=1):
            title = normalized_text(slide.shapes.title.text) if slide.shapes.title else ""
            title = title or f"Slide {slide_number}"
            detected_title = detected_title or (title if not title.startswith("Slide ") else None)
            heading_path = (title,)
            location = SourceLocation(section_path=heading_path, slide_number=slide_number)
            blocks.append(
                StructuredBlock(
                    block_id=f"block_{len(blocks) + 1:04d}",
                    kind=BlockKind.HEADING,
                    text=title,
                    heading_path=heading_path,
                    location=location,
                    metadata={"level": 1},
                )
            )
            indexed_shapes = list(enumerate(slide.shapes))
            ordered_shapes = sorted(
                indexed_shapes,
                key=lambda item: (
                    int(getattr(item[1], "top", 0) or 0),
                    int(getattr(item[1], "left", 0) or 0),
                    item[0],
                ),
            )
            for _, shape in ordered_shapes:
                if slide.shapes.title is shape:
                    continue
                shape_count += self._append_shape(
                    shape,
                    slide_number,
                    heading_path,
                    blocks,
                    images,
                    tables,
                    image_by_hash,
                )
                if shape_count > self.limits.max_shapes:
                    self._raise_limit()
            if slide.has_notes_slide:
                note_parts = [
                    normalized_text(shape.text)
                    for shape in slide.notes_slide.shapes
                    if getattr(shape, "has_text_frame", False) and normalized_text(shape.text)
                ]
                notes = "\n".join(dict.fromkeys(note_parts))
                if notes:
                    blocks.append(
                        StructuredBlock(
                            block_id=f"block_{len(blocks) + 1:04d}",
                            kind=BlockKind.PARAGRAPH,
                            text=notes,
                            heading_path=heading_path,
                            location=location,
                            metadata={"speaker_notes": True},
                        )
                    )

        return complete_document(
            request=request,
            source_format=self.source_format,
            parser_name=self.parser_id,
            parser_version=self.parser_version,
            provider_name=request.route.provider,
            provider_version="python-pptx-1.0",
            blocks=blocks,
            images=images,
            tables=tables,
            detected_title=detected_title or request.filename.rsplit(".", 1)[0],
            slides=len(presentation.slides),
            started_at=started,
            reference_knowhere=True,
        )

    def _append_shape(
        self,
        shape,
        slide_number: int,
        heading_path: tuple[str, ...],
        blocks: list[StructuredBlock],
        images: list[ImageAssetDraft],
        tables: list[TableAssetDraft],
        image_by_hash: dict[str, ImageAssetDraft],
    ) -> int:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            count = 1
            for child in shape.shapes:
                count += self._append_shape(
                    child,
                    slide_number,
                    heading_path,
                    blocks,
                    images,
                    tables,
                    image_by_hash,
                )
            return count
        location = SourceLocation(
            section_path=heading_path,
            slide_number=slide_number,
            bbox=self._bbox(shape),
        )
        if getattr(shape, "has_table", False):
            rows = [[normalized_text(cell.text) for cell in row.cells] for row in shape.table.rows]
            markdown, headers = matrix_to_markdown(rows)
            if markdown:
                asset_id = f"table_{len(tables) + 1:04d}"
                table = TableAssetDraft(
                    asset_id=asset_id,
                    title=f"{heading_path[0]} / Table {len(tables) + 1}",
                    html=matrix_to_html(rows),
                    markdown=markdown,
                    headers=headers,
                    row_count=max(len(rows) - 1, 0),
                    column_count=max((len(row) for row in rows), default=0),
                    location=location,
                )
                tables.append(table)
                blocks.append(
                    StructuredBlock(
                        block_id=f"block_{len(blocks) + 1:04d}",
                        kind=BlockKind.TABLE,
                        text=markdown,
                        heading_path=heading_path,
                        location=location,
                        table_asset_ids=(asset_id,),
                    )
                )
            return 1
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            asset = self._picture_asset(shape, heading_path, location, images, image_by_hash)
            blocks.append(
                StructuredBlock(
                    block_id=f"block_{len(blocks) + 1:04d}",
                    kind=BlockKind.IMAGE,
                    text=asset.caption or asset.payload.filename,
                    heading_path=heading_path,
                    location=location,
                    image_asset_ids=(asset.asset_id,),
                )
            )
            return 1
        if getattr(shape, "has_text_frame", False):
            text = normalized_text(shape.text)
            if text:
                blocks.append(
                    StructuredBlock(
                        block_id=f"block_{len(blocks) + 1:04d}",
                        kind=BlockKind.PARAGRAPH,
                        text=text,
                        heading_path=heading_path,
                        location=location,
                        metadata={"shape_name": normalized_text(shape.name)},
                    )
                )
        return 1

    def _picture_asset(
        self,
        shape,
        heading_path: tuple[str, ...],
        location: SourceLocation,
        images: list[ImageAssetDraft],
        image_by_hash: dict[str, ImageAssetDraft],
    ) -> ImageAssetDraft:
        content = bytes(shape.image.blob)
        digest = hashlib.sha256(content).hexdigest()
        if digest in image_by_hash:
            return image_by_hash[digest]
        if len(images) >= self.limits.max_images:
            self._raise_limit()
        inspection = inspect_image(content)
        filename = f"slide-{location.slide_number}-image-{len(images) + 1}.{shape.image.ext}"
        payload = self.payload_store.put(filename, inspection.content_type, content)
        properties = shape._element.xpath(".//p:cNvPr")
        alt_text = ""
        if properties:
            alt_text = normalized_text(properties[0].get("descr") or properties[0].get("title"))
        caption = alt_text or normalized_text(shape.name) or " / ".join(heading_path)
        asset = ImageAssetDraft(
            asset_id=f"image_{len(images) + 1:04d}",
            payload=payload,
            image_type=f"pptx_embedded:{inspection.width}x{inspection.height}",
            caption=caption,
            caption_source=CaptionSource.PARSER,
            caption_confidence=0.9 if alt_text else 0.55,
            location=location,
        )
        images.append(asset)
        image_by_hash[digest] = asset
        return asset

    @staticmethod
    def _bbox(shape) -> tuple[float, float, float, float]:
        return tuple(
            round(float(value or 0) / _EMU_PER_POINT, 2)
            for value in (shape.left, shape.top, shape.width, shape.height)
        )

    @staticmethod
    def _raise_limit() -> None:
        raise IngestError(
            IngestErrorCode.DOCUMENT_LIMIT_EXCEEDED,
            "The PPTX exceeds configured slide, shape, or image limits.",
        )
