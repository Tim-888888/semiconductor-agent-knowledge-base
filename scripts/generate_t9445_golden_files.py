"""Generate deterministic, publicly displayable T9-4.4.5 golden documents."""

from __future__ import annotations

import argparse
import hashlib
import io
import json
import re
import zipfile
import zlib
from datetime import UTC, datetime
from pathlib import Path

from docx import Document
from docx.enum.text import WD_BREAK
from docx.shared import Inches
from openpyxl import Workbook
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches as PptxInches

GOLDEN_VERSION = "t9-4.4.5-golden-v1"
FIXED_TIME = datetime(2026, 1, 1, tzinfo=UTC)
ZIP_TIME = (2026, 1, 1, 0, 0, 0)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--output-dir",
        type=Path,
        default=Path("data/t9445_golden"),
    )
    return parser.parse_args()


def _sha256(content: bytes) -> str:
    return hashlib.sha256(content).hexdigest()


def _normalize_ooxml(content: bytes) -> bytes:
    """Remove ZIP timestamp entropy without altering package members."""

    source = io.BytesIO(content)
    output = io.BytesIO()
    with zipfile.ZipFile(source) as archive, zipfile.ZipFile(
        output,
        "w",
        compression=zipfile.ZIP_DEFLATED,
        compresslevel=9,
    ) as normalized:
        for name in sorted(archive.namelist()):
            original = archive.getinfo(name)
            item = zipfile.ZipInfo(name, ZIP_TIME)
            item.compress_type = zipfile.ZIP_DEFLATED
            item.external_attr = original.external_attr
            item.create_system = original.create_system
            payload = archive.read(name)
            if name == "docProps/core.xml":
                payload = re.sub(
                    rb"\d{4}-\d{2}-\d{2}T\d{2}:\d{2}:\d{2}Z",
                    b"2026-01-01T00:00:00Z",
                    payload,
                )
            normalized.writestr(item, payload)
    return output.getvalue()


def _wafer_png() -> bytes:
    image = Image.new("RGB", (960, 720), "white")
    draw = ImageDraw.Draw(image)
    draw.text((30, 24), "T9445 IMAGE EDGE RING 57", fill="#17212b")
    draw.text((30, 50), "ETCH-03 Chamber B wafer inspection map", fill="#17212b")
    center_x, center_y, radius = 480, 360, 260
    draw.ellipse(
        (center_x - radius, center_y - radius, center_x + radius, center_y + radius),
        outline="#243746",
        width=5,
    )
    for row in range(-10, 11):
        for column in range(-10, 11):
            x = center_x + column * 22
            y = center_y + row * 22
            distance = (column**2 + row**2) ** 0.5
            if distance > 10.2:
                continue
            color = "#d94b40" if distance >= 8.2 else "#4f9d85"
            draw.rectangle((x - 8, y - 8, x + 8, y + 8), fill=color)
    draw.text((30, 680), "Finding: edge-ring concentration; center region stable", fill="#17212b")
    output = io.BytesIO()
    image.save(output, format="PNG", optimize=False)
    return output.getvalue()


def _docx(wafer_png: bytes) -> bytes:
    document = Document()
    document.core_properties.title = "T9445 DOCX Chamber Recovery SOP"
    document.core_properties.created = FIXED_TIME
    document.core_properties.modified = FIXED_TIME
    document.add_heading("ETCH-03 Chamber B Recovery", level=1)
    document.add_paragraph(
        "Control token T9445-DOCX-LEAKCHECK-33. Hold the next lot and verify chamber "
        "pressure before recipe recovery."
    )
    document.add_paragraph("Complete a leak check.", style="List Bullet")
    document.add_paragraph("Review RF match trend.", style="List Bullet")
    table = document.add_table(rows=3, cols=3)
    for column, value in enumerate(("Signal", "Limit", "Action")):
        table.cell(0, column).text = value
    for column, value in enumerate(("Pressure", "12 Pa", "Hold lot")):
        table.cell(1, column).text = value
    for column, value in enumerate(("RF Match", "85 %", "Escalate")):
        table.cell(2, column).text = value
    run = document.add_paragraph("Controlled wafer evidence:").add_run()
    inline = run.add_picture(io.BytesIO(wafer_png), width=Inches(3.2))
    inline._inline.docPr.set("descr", "T9445 DOCX edge-ring wafer evidence")
    document.add_paragraph().add_run().add_break(WD_BREAK.PAGE)
    document.add_heading("Release Criteria", level=2)
    document.add_paragraph("Release only after pressure and RF match return inside limits.")
    output = io.BytesIO()
    document.save(output)
    return _normalize_ooxml(output.getvalue())


def _xlsx() -> bytes:
    workbook = Workbook()
    workbook.properties.title = "T9445 XLSX FDC limits"
    workbook.properties.created = FIXED_TIME
    workbook.properties.modified = FIXED_TIME
    fdc = workbook.active
    fdc.title = "FDC_Limits"
    fdc.append(["Signal", "Limit(Pa)", "RF Match(%)", "Control Token"])
    fdc.append(["Chamber Pressure", 12, 85, "T9445-XLSX-PRESSURE-64"])
    fdc.append(["Pressure Recovery", "=B2-1", 92, "Review after clean"])
    fdc.merge_cells("A5:D5")
    fdc["A5"] = "Escalation Matrix"
    fdc.append(["Condition", "Owner", "Deadline(min)", "Disposition"])
    fdc.append(["Edge ring", "Process", 15, "Hold lot"])
    recipe = workbook.create_sheet("Recipe_Audit")
    recipe.append(["Recipe", "Version", "Approved"])
    recipe.append(["ETCH-ALPHA", "V2.3", "yes"])
    hidden = workbook.create_sheet("Internal_Raw")
    hidden.sheet_state = "hidden"
    hidden["A1"] = "not for retrieval"
    output = io.BytesIO()
    workbook.save(output)
    workbook.close()
    return _normalize_ooxml(output.getvalue())


def _pptx(wafer_png: bytes) -> bytes:
    presentation = Presentation()
    presentation.core_properties.title = "T9445 PPTX Edge Ring Training"
    presentation.core_properties.created = FIXED_TIME
    presentation.core_properties.modified = FIXED_TIME
    first = presentation.slides.add_slide(presentation.slide_layouts[5])
    first.shapes.title.text = "ETCH-03 Edge-Ring Response"
    box = first.shapes.add_textbox(PptxInches(0.7), PptxInches(1.2), PptxInches(7), PptxInches(0.8))
    box.text_frame.text = (
        "Control token T9445-PPTX-EDGE-72. Hold the lot and correlate pressure with RF match."
    )
    table = first.shapes.add_table(
        3,
        2,
        PptxInches(0.7),
        PptxInches(2.2),
        PptxInches(4.5),
        PptxInches(1.6),
    ).table
    for column, value in enumerate(("Check", "Expected")):
        table.cell(0, column).text = value
    for row, values in enumerate((("Pressure", "<= 12 Pa"), ("RF Match", ">= 85 %")), start=1):
        for column, value in enumerate(values):
            table.cell(row, column).text = value
    picture = first.shapes.add_picture(
        io.BytesIO(wafer_png),
        PptxInches(5.5),
        PptxInches(2.0),
        width=PptxInches(3.8),
    )
    picture._element.nvPicPr.cNvPr.set("descr", "T9445 PPTX edge-ring wafer map")
    second = presentation.slides.add_slide(presentation.slide_layouts[5])
    second.shapes.title.text = "Release Gate"
    gate = second.shapes.add_textbox(
        PptxInches(0.7), PptxInches(1.3), PptxInches(8), PptxInches(1.0)
    )
    gate.text_frame.text = "Release after leak check completion and two stable monitor wafers."
    output = io.BytesIO()
    presentation.save(output)
    return _normalize_ooxml(output.getvalue())


def _pdf_stream(payload: bytes) -> bytes:
    return b"<< /Length " + str(len(payload)).encode() + b" >>\nstream\n" + payload + b"\nendstream"


def _pdf(wafer_png: bytes) -> bytes:
    image = Image.open(io.BytesIO(wafer_png)).convert("RGB").resize((240, 180))
    image_payload = zlib.compress(image.tobytes(), level=9)
    page_one = b"\n".join(
        (
            b"BT /F1 18 Tf 60 740 Td (T9445 PDF Controlled Case) Tj ET",
            b"BT /F1 11 Tf 60 710 Td (Control token T9445-PDF-ESCALATE-42) Tj ET",
            b"BT /F1 11 Tf 60 686 Td (ETCH-03 Chamber B pressure alarm after chamber clean.) Tj ET",
            b"BT /F1 11 Tf 60 662 Td (Continue on page 2 for the release gate and wafer evidence.) Tj ET",
            b"BT /F1 11 Tf 60 638 Td (Leak check is mandatory before release.) Tj ET",
        )
    )
    page_two = b"\n".join(
        (
            b"BT /F1 16 Tf 60 740 Td (Release Gate and Evidence) Tj ET",
            b"BT /F1 10 Tf 60 700 Td (Signal) Tj 150 0 Td (Limit) Tj 120 0 Td (Action) Tj ET",
            b"BT /F1 10 Tf 60 675 Td (Pressure) Tj 150 0 Td (12 Pa) Tj 120 0 Td (Hold lot) Tj ET",
            b"BT /F1 10 Tf 60 650 Td (RF Match) Tj 150 0 Td (85 percent) Tj 120 0 Td (Escalate) Tj ET",
            b"0.5 w 55 635 410 90 re S 55 665 m 465 665 l S 200 635 m 200 725 l S 320 635 m 320 725 l S",
            b"q 240 0 0 180 300 390 cm /Im1 Do Q",
            b"BT /F1 10 Tf 60 420 Td (Embedded wafer map: edge-ring defect concentration.) Tj ET",
            b"BT /F1 11 Tf 60 380 Td (Release only after leak check and two stable monitor wafers.) Tj ET",
        )
    )
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R 4 0 R] /Count 2 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 5 0 R >> >> /Contents 7 0 R >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 5 0 R >> /XObject << /Im1 6 0 R >> >> /Contents 8 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        (
            b"<< /Type /XObject /Subtype /Image /Width 240 /Height 180 "
            b"/ColorSpace /DeviceRGB /BitsPerComponent 8 /Filter /FlateDecode /Length "
            + str(len(image_payload)).encode()
            + b" >>\nstream\n"
            + image_payload
            + b"\nendstream"
        ),
        _pdf_stream(page_one),
        _pdf_stream(page_two),
    ]
    payload = bytearray(b"%PDF-1.7\n%\xe2\xe3\xcf\xd3\n")
    offsets = [0]
    for index, body in enumerate(objects, start=1):
        offsets.append(len(payload))
        payload.extend(f"{index} 0 obj\n".encode())
        payload.extend(body)
        payload.extend(b"\nendobj\n")
    xref_offset = len(payload)
    payload.extend(f"xref\n0 {len(objects) + 1}\n".encode())
    payload.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        payload.extend(f"{offset:010d} 00000 n \n".encode())
    payload.extend(
        (
            f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
            f"startxref\n{xref_offset}\n%%EOF\n"
        ).encode()
    )
    return bytes(payload)


def _entries(wafer_png: bytes) -> list[dict[str, object]]:
    markdown = b"""# T9445 Markdown Pressure SOP

Control token `T9445-MD-PRESSURE-17`.

## First Wafer Gate

- Hold the lot after a pressure alarm.
- Review RF match before release.

| Signal | Limit | Unit |
| --- | ---: | --- |
| Chamber Pressure | 12 | Pa |

```text
release = leak_check_passed and monitor_wafers >= 2
```
"""
    text = (
        b"T9445 TXT Chamber Checklist\n\n"
        b"Control token T9445-TXT-RFMATCH-82.\n"
        b"After chamber clean, confirm RF match is at least 85 percent.\n"
        b"Hold the lot until two monitor wafers are stable.\n"
    )
    html = b"""<!doctype html><html><head><title>T9445 HTML Alarm Guide</title>
<script>navigationNoise()</script></head><body><nav>uncontrolled menu</nav><main>
<h1>Pressure Alarm Response</h1><p>Control token T9445-HTML-FORELINE-29.</p>
<h2>Checks</h2><p>Inspect foreline pressure and chamber leak status.</p>
<table><tr><th>Signal</th><th>Limit</th></tr><tr><td>Foreline</td><td>9 Pa</td></tr></table>
</main></body></html>"""
    csv_content = (
        b"Wafer;Pressure(Pa);RF Match(%);Disposition;Control Token\n"
        b"W01;10.8;93;PASS;T9445-CSV-SPC-58\n"
        b"W02;12.7;81;HOLD;edge-ring-review\n"
    )
    return [
        {
            "category": "markdown",
            "filename": "01_pressure_sop.md",
            "content_type": "text/markdown",
            "document_id": "T9445-GOLDEN-MD",
            "title": "T9-4.4.5 Markdown Pressure SOP",
            "document_type": "sop",
            "parser_name": "markdown-structured-v1",
            "content": markdown,
            "required_terms": ["T9445-MD-PRESSURE-17", "First Wafer Gate", "Chamber Pressure"],
            "min_chunks": 3,
            "min_tables": 1,
            "min_images": 0,
            "retrieval_query": "T9445-MD-PRESSURE-17 first wafer pressure gate",
        },
        {
            "category": "text",
            "filename": "02_rf_match_checklist.txt",
            "content_type": "text/plain",
            "document_id": "T9445-GOLDEN-TXT",
            "title": "T9-4.4.5 RF Match Checklist",
            "document_type": "checklist",
            "parser_name": "text-structured-v1",
            "content": text,
            "required_terms": ["T9445-TXT-RFMATCH-82", "RF match", "monitor wafers"],
            "min_chunks": 1,
            "min_tables": 0,
            "min_images": 0,
            "retrieval_query": "T9445-TXT-RFMATCH-82 RF match monitor wafers",
        },
        {
            "category": "html",
            "filename": "03_foreline_alarm.html",
            "content_type": "text/html",
            "document_id": "T9445-GOLDEN-HTML",
            "title": "T9-4.4.5 HTML Foreline Alarm Guide",
            "document_type": "training_note",
            "parser_name": "html-structured-v1",
            "content": html,
            "required_terms": ["T9445-HTML-FORELINE-29", "foreline pressure", "9 Pa"],
            "forbidden_terms": ["navigationNoise", "uncontrolled menu"],
            "min_chunks": 2,
            "min_tables": 1,
            "min_images": 0,
            "retrieval_query": "T9445-HTML-FORELINE-29 foreline pressure limit",
        },
        {
            "category": "pdf",
            "filename": "04_controlled_case.pdf",
            "content_type": "application/pdf",
            "document_id": "T9445-GOLDEN-PDF-V2",
            "title": "T9-4.4.5 PDF Controlled Case",
            "document_type": "case",
            "parser_name": "pdf-mineru-v1",
            "content": _pdf(wafer_png),
            "required_terms": ["T9445-PDF-ESCALATE-42", "Release Gate", "leak check"],
            "min_chunks": 2,
            "min_tables": 1,
            "min_images": 1,
            "retrieval_query": "T9445-PDF-ESCALATE-42 release gate leak check",
        },
        {
            "category": "docx",
            "filename": "05_chamber_recovery.docx",
            "content_type": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "document_id": "T9445-GOLDEN-DOCX",
            "title": "T9-4.4.5 DOCX Chamber Recovery SOP",
            "document_type": "sop",
            "parser_name": "docx-structured-v1",
            "content": _docx(wafer_png),
            "required_terms": ["T9445-DOCX-LEAKCHECK-33", "Recovery", "12 Pa"],
            "min_chunks": 5,
            "min_tables": 1,
            "min_images": 1,
            "retrieval_query": "T9445-DOCX-LEAKCHECK-33 chamber recovery leak check",
        },
        {
            "category": "xlsx",
            "filename": "06_fdc_limits.xlsx",
            "content_type": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "document_id": "T9445-GOLDEN-XLSX",
            "title": "T9-4.4.5 XLSX FDC Limits",
            "document_type": "fdc_limits",
            "parser_name": "xlsx-structured-v1",
            "content": _xlsx(),
            "required_terms": ["T9445-XLSX-PRESSURE-64", "FDC_Limits", "ETCH-ALPHA"],
            "min_chunks": 3,
            "min_tables": 3,
            "min_images": 0,
            "retrieval_query": "T9445-XLSX-PRESSURE-64 FDC chamber pressure limit",
        },
        {
            "category": "csv",
            "filename": "07_spc_samples.csv",
            "content_type": "text/csv",
            "document_id": "T9445-GOLDEN-CSV",
            "title": "T9-4.4.5 CSV SPC Samples",
            "document_type": "spc_data",
            "parser_name": "csv-structured-v1",
            "content": csv_content,
            "required_terms": ["T9445-CSV-SPC-58", "Pressure(Pa)", "HOLD"],
            "min_chunks": 1,
            "min_tables": 1,
            "min_images": 0,
            "retrieval_query": "T9445-CSV-SPC-58 SPC pressure disposition",
        },
        {
            "category": "pptx",
            "filename": "08_edge_ring_training.pptx",
            "content_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "document_id": "T9445-GOLDEN-PPTX",
            "title": "T9-4.4.5 PPTX Edge Ring Training",
            "document_type": "training_note",
            "parser_name": "pptx-structured-v1",
            "content": _pptx(wafer_png),
            "required_terms": ["T9445-PPTX-EDGE-72", "Release Gate", "RF Match"],
            "min_chunks": 4,
            "min_tables": 1,
            "min_images": 1,
            "retrieval_query": "T9445-PPTX-EDGE-72 edge ring release gate",
        },
        {
            "category": "image",
            "filename": "09_edge_ring_wafer.png",
            "content_type": "image/png",
            "document_id": "T9445-GOLDEN-IMAGE",
            "title": "T9-4.4.5 Edge Ring Wafer Map",
            "document_type": "wafer_map",
            "parser_name": "image-vlm-v1",
            "content": wafer_png,
            "required_terms": ["edge", "ring"],
            "min_chunks": 1,
            "min_tables": 0,
            "min_images": 1,
            "retrieval_query": "T9445 IMAGE EDGE RING 57 ETCH-03 wafer map",
        },
    ]


def generate(output_dir: Path) -> dict[str, object]:
    output_dir.mkdir(parents=True, exist_ok=True)
    wafer_png = _wafer_png()
    manifest_entries: list[dict[str, object]] = []
    for entry in _entries(wafer_png):
        content = entry.pop("content")
        assert isinstance(content, bytes)
        path = output_dir / str(entry["filename"])
        path.write_bytes(content)
        manifest_entries.append(
            {
                **entry,
                "revision": "R1",
                "source_kind": "synthetic_acceptance",
                "source_license": "CC0-1.0",
                "source_uri": f"synthetic://{GOLDEN_VERSION}/{entry['filename']}",
                "sha256": _sha256(content),
                "size_bytes": len(content),
            }
        )
    manifest = {
        "schema_version": 1,
        "golden_version": GOLDEN_VERSION,
        "generated_at": FIXED_TIME.isoformat(),
        "license": "CC0-1.0",
        "scope": "Synthetic semiconductor multi-format ingestion acceptance data.",
        "entries": manifest_entries,
    }
    (output_dir / "manifest.json").write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )
    return manifest


def main() -> None:
    args = parse_args()
    manifest = generate(args.output_dir)
    print(
        json.dumps(
            {
                "output_dir": str(args.output_dir.resolve()),
                "golden_version": manifest["golden_version"],
                "files": len(manifest["entries"]),
            },
            ensure_ascii=False,
        )
    )


if __name__ == "__main__":
    main()
