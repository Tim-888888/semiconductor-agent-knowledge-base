from __future__ import annotations

import io
import json
from dataclasses import dataclass

import pytest
from docx import Document as WordDocument
from openpyxl import Workbook
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from semikb.config import Settings
from semikb.contracts.models import DocumentLifecycle, IngestionStatus
from semikb.rag_ingestion.semikb_adapter import SemikbIngestAdapter
from semikb.rag_ingestion.service import IngestionService
from semikb.rag_retrieval.encoders import DeterministicHybridEncoder
from semikb.storage.memory import DemoStore
from semikb_ingest import IngestError, IngestErrorCode
from semikb_ingest.providers import (
    MinerUContentItem,
    MinerUImage,
    MinerUPdfResult,
    ProviderRegistry,
    VisionAnalysis,
)


def _png_bytes() -> bytes:
    output = io.BytesIO()
    Image.new("RGB", (32, 32), (218, 54, 46)).save(output, format="PNG")
    return output.getvalue()


def _docx_bytes() -> bytes:
    document = WordDocument()
    document.core_properties.title = "ETCH chamber recovery"
    document.add_heading("Recovery SOP", level=1)
    document.add_paragraph("Verify chamber pressure before releasing the next lot.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Signal"
    table.cell(0, 1).text = "Limit"
    table.cell(1, 0).text = "Pressure"
    table.cell(1, 1).text = "12"
    document.add_picture(io.BytesIO(_png_bytes()), width=Inches(0.5))
    output = io.BytesIO()
    document.save(output)
    return output.getvalue()


def _xlsx_bytes() -> bytes:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "FDC"
    sheet.append(["Signal", "Limit"])
    sheet.append(["Pressure", 12])
    output = io.BytesIO()
    workbook.save(output)
    workbook.close()
    return output.getvalue()


def _pptx_bytes() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "ETCH alarm response"
    box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(6), Inches(1))
    box.text_frame.text = "Check pressure and RF match before recipe recovery."
    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue()


@dataclass
class FakeVisionProvider:
    provider_name: str = "qwen-vl"
    provider_version: str = "qwen3.7-plus"

    def analyze_image(self, **_kwargs) -> VisionAnalysis:
        return VisionAnalysis(
            caption="晶圆边缘出现环状缺陷。",
            ocr_text="ETCH-03",
            detection_summary="边缘缺陷集中。",
            confidence=0.93,
            detected_language="zh",
        )


@dataclass
class FakePdfProvider:
    provider_name: str = "mineru"
    provider_version: str = "vlm"

    def parse_pdf(self, **_kwargs) -> MinerUPdfResult:
        return MinerUPdfResult(
            markdown="# ETCH SOP\n\nPressure alarm response.\n",
            pages=2,
            content_items=(
                MinerUContentItem(
                    kind="text",
                    text="ETCH SOP",
                    heading_level=1,
                    page_number=1,
                ),
                MinerUContentItem(
                    kind="text",
                    text="Pressure alarm response.",
                    page_number=1,
                ),
                MinerUContentItem(
                    kind="table",
                    table_html=(
                        "<table><tr><th>Signal</th><th>Limit</th></tr>"
                        "<tr><td>Pressure</td><td>12</td></tr></table>"
                    ),
                    table_caption="Alarm limits",
                    page_number=2,
                ),
                MinerUContentItem(
                    kind="image",
                    image_path="images/wafer.png",
                    page_number=2,
                ),
            ),
            images=(
                MinerUImage(
                    path="images/wafer.png",
                    filename="wafer.png",
                    content_type="image/png",
                    content=_png_bytes(),
                    caption="Wafer edge map",
                    page_number=2,
                ),
            ),
        )


@dataclass
class TimeoutPdfProvider:
    provider_name: str = "mineru"
    provider_version: str = "vlm"

    def parse_pdf(self, **_kwargs) -> MinerUPdfResult:
        raise IngestError(
            IngestErrorCode.PARSER_TIMEOUT,
            "The document parser timed out.",
        )


def _registry(*, timeout_pdf: bool = False) -> ProviderRegistry:
    registry = ProviderRegistry()
    registry.register(TimeoutPdfProvider() if timeout_pdf else FakePdfProvider())
    registry.register(FakeVisionProvider())
    return registry


def _service(
    store: DemoStore,
    *,
    registry: ProviderRegistry | None = None,
) -> IngestionService:
    settings = Settings(_env_file=None, demo_mode=True, embedding_dim=8)
    return IngestionService(
        store,
        settings,
        encoder=DeterministicHybridEncoder(8),
        ingest_adapter=SemikbIngestAdapter(settings, registry),
    )


def _metadata(document_id: str) -> dict[str, str]:
    return {
        "document_id": document_id,
        "revision": "R1",
        "title": "T9-4.4.4 integration fixture",
        "document_type": "sop",
        "approval_status": "approved",
        "lifecycle": "published",
        "access_scope_key": "demo_engineering",
        "fab": "FAB-01",
        "product": "P-ALPHA",
        "tool_id": "ETCH-03",
    }


@pytest.mark.parametrize(
    ("filename", "content", "parser_name"),
    [
        ("note.md", b"# Alarm\n\nCheck chamber pressure.", "markdown-structured-v1"),
        ("note.txt", b"Alarm response\nCheck chamber pressure.", "text-structured-v1"),
        (
            "note.html",
            b"<h1>Alarm</h1><p>Check chamber pressure.</p>",
            "html-structured-v1",
        ),
        ("limits.csv", b"Signal,Limit\nPressure,12\n", "csv-structured-v1"),
        ("sop.docx", _docx_bytes(), "docx-structured-v1"),
        ("fdc.xlsx", _xlsx_bytes(), "xlsx-structured-v1"),
        ("training.pptx", _pptx_bytes(), "pptx-structured-v1"),
    ],
    ids=["markdown", "text", "html", "csv", "docx", "xlsx", "pptx"],
)
def test_local_format_outputs_flow_through_governed_publication(
    filename: str,
    content: bytes,
    parser_name: str,
) -> None:
    store = DemoStore()
    document_id = f"T9444-{parser_name.split('-', 1)[0].upper()}"

    job = _service(store).ingest_file(filename, content, _metadata(document_id))

    document = store.get_document(document_id, "R1")
    chunks = [item for item in store.chunks.values() if item.document_id == document_id]
    assert job.status is IngestionStatus.PUBLISHED
    assert job.parse_contract_version == "semikb-ingest-v1"
    assert job.parser_name == parser_name
    assert job.parser_version != "pending"
    assert job.chunker_version == "structured-blocks-v1"
    assert job.parse_metrics["chunks"] == job.chunks_count
    assert document is not None
    assert document.lifecycle is DocumentLifecycle.PUBLISHED
    assert document.parser_name == parser_name
    assert document.parsed_ref is not None
    assert store.load_object(document.parsed_ref)
    assert chunks
    assert all(item.parser_name == parser_name for item in chunks)
    assert all(item.lifecycle is DocumentLifecycle.PUBLISHED for item in chunks)
    if parser_name in {
        "markdown-structured-v1",
        "docx-structured-v1",
        "xlsx-structured-v1",
        "pptx-structured-v1",
    }:
        assert job.upstream_commit == "2e4eb5846249d273b11902ee00f26db949e45b38"
        assert document.upstream_commit == job.upstream_commit


def test_table_and_image_assets_keep_stable_links_and_durable_objects() -> None:
    store = DemoStore()
    service = _service(store)

    job = service.ingest_file("sop.docx", _docx_bytes(), _metadata("T9444-ASSETS"))

    tables = [item for item in store.tables.values() if item.document_id == "T9444-ASSETS"]
    images = [item for item in store.images.values() if item.document_id == "T9444-ASSETS"]
    chunks = [item for item in store.chunks.values() if item.document_id == "T9444-ASSETS"]
    assert job.status is IngestionStatus.PUBLISHED
    assert job.tables_count == len(tables) == 1
    assert job.images_count == len(images) == 1
    assert tables[0].table_id in {item for chunk in chunks for item in chunk.table_ids}
    assert images[0].image_id in {item for chunk in chunks for item in chunk.image_ids}
    table_payload = json.loads(store.load_object(tables[0].object_ref))
    assert table_payload["table_id"] == tables[0].table_id
    assert table_payload["headers"] == ["Signal", "Limit"]
    assert store.load_object(images[0].object_ref) == _png_bytes()
    assert tables[0].lifecycle is DocumentLifecycle.PUBLISHED
    assert images[0].lifecycle is DocumentLifecycle.PUBLISHED


@pytest.mark.parametrize(
    ("filename", "content", "document_id", "parser_name"),
    [
        ("sop.pdf", b"%PDF-1.7\nsynthetic", "T9444-PDF", "pdf-mineru-v1"),
        ("wafer.png", _png_bytes(), "T9444-IMAGE", "image-vlm-v1"),
    ],
)
def test_provider_formats_publish_with_provider_audit_and_assets(
    filename: str,
    content: bytes,
    document_id: str,
    parser_name: str,
) -> None:
    store = DemoStore()

    job = _service(store, registry=_registry()).ingest_file(
        filename,
        content,
        _metadata(document_id),
    )

    document = store.get_document(document_id, "R1")
    assert job.status is IngestionStatus.PUBLISHED
    assert job.parser_name == parser_name
    assert job.provider_name in {"mineru", "qwen-vl"}
    assert job.provider_version in {"vlm", "qwen3.7-plus"}
    assert job.images_count >= 1
    assert document is not None
    assert document.provider_name == job.provider_name
    assert document.provider_version == job.provider_version
    if parser_name == "image-vlm-v1":
        assert job.upstream_commit == "2e4eb5846249d273b11902ee00f26db949e45b38"
    assert all(
        image.lifecycle is DocumentLifecycle.PUBLISHED
        for image in store.images.values()
        if image.document_id == document_id
    )


def test_parser_timeout_remains_unpublished_and_retry_is_idempotent() -> None:
    store = DemoStore()
    content = b"%PDF-1.7\nsynthetic"
    metadata = _metadata("T9444-RETRY")
    failed = _service(store, registry=_registry(timeout_pdf=True)).ingest_file(
        "sop.pdf",
        content,
        metadata,
    )

    assert failed.status is IngestionStatus.FAILED
    assert failed.error_code == IngestErrorCode.PARSER_TIMEOUT.value
    assert store.get_document("T9444-RETRY", "R1") is None
    assert not [item for item in store.chunks.values() if item.document_id == "T9444-RETRY"]
    assert not [item for item in store.images.values() if item.document_id == "T9444-RETRY"]
    assert not [item for item in store.tables.values() if item.document_id == "T9444-RETRY"]

    service = _service(store, registry=_registry())
    completed = service.retry(failed.job_id)
    repeated = service.ingest_file("sop.pdf", content, metadata)

    assert completed.status is IngestionStatus.PUBLISHED
    assert repeated.job_id == completed.job_id
    assert repeated.attempt == 2
    assert len(
        [item for item in store.chunks.values() if item.document_id == "T9444-RETRY"]
    ) == completed.chunks_count
    assert len(
        [item for item in store.images.values() if item.document_id == "T9444-RETRY"]
    ) == completed.images_count
    assert len(
        [item for item in store.tables.values() if item.document_id == "T9444-RETRY"]
    ) == completed.tables_count
