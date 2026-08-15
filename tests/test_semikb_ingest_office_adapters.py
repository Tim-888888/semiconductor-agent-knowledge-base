from __future__ import annotations

import io
import json
from pathlib import Path

from docx import Document
from jsonschema import Draft202012Validator
from openpyxl import Workbook
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from semikb_ingest import ChunkType, SourceFormat, build_dispatcher
from semikb_ingest.assets import ProcessPayloadStore

_CONTRACT = json.loads(
    Path("docs/evidence/t9-4-4-1/semikb-ingest-contract-v1.schema.json").read_text(encoding="utf-8")
)


def _assert_frozen_contract(parsed) -> None:
    Draft202012Validator(_CONTRACT).validate(parsed.model_dump(mode="json"))


def _png_bytes(color: tuple[int, int, int] = (220, 40, 40)) -> bytes:
    output = io.BytesIO()
    Image.new("RGB", (48, 32), color).save(output, format="PNG")
    return output.getvalue()


def _docx_bytes() -> bytes:
    document = Document()
    document.core_properties.title = "ETCH-03 Maintenance"
    document.add_heading("Chamber Clean", level=1)
    document.add_paragraph("Verify pressure and RF match.")
    document.add_paragraph("Hold the next lot.", style="List Bullet")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Item"
    table.cell(0, 1).text = "Limit"
    table.cell(1, 0).text = "Pressure"
    table.cell(1, 1).text = "12 Pa"
    run = document.add_paragraph("Wafer map:").add_run()
    run.add_picture(io.BytesIO(_png_bytes()), width=Inches(1))
    output = io.BytesIO()
    document.save(output)
    return output.getvalue()


def _xlsx_bytes() -> bytes:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "FDC"
    sheet.append(["Pressure(kPa)", "RF Match(%)"])
    sheet.append([10.2, 98])
    sheet.append([12.1, "=B2-20"])
    sheet.merge_cells("A5:B5")
    sheet["A5"] = "Second Region"
    sheet.append(["Signal", "State"])
    sheet.append(["Pressure", "Alarm"])
    hidden = workbook.create_sheet("Internal")
    hidden.sheet_state = "hidden"
    hidden["A1"] = "must not index"
    output = io.BytesIO()
    workbook.save(output)
    return output.getvalue()


def _pptx_bytes() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Wafer Edge Defect"
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4), Inches(0.8))
    textbox.text = "Inspect edge ring and chamber pressure."
    table = slide.shapes.add_table(2, 2, Inches(0.5), Inches(2.2), Inches(4), Inches(1.2)).table
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Edge defects"
    table.cell(1, 1).text = "18"
    slide.shapes.add_picture(
        io.BytesIO(_png_bytes((40, 160, 80))), Inches(5), Inches(1.2), width=Inches(1)
    )
    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue()


def test_docx_preserves_document_order_headings_tables_and_images() -> None:
    payload_store = ProcessPayloadStore()
    parsed = build_dispatcher(payload_store).parse(
        "maintenance.docx",
        _docx_bytes(),
        declared_media_type=(
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        ),
        correlation_id="docx-1",
    )

    assert parsed.source_format is SourceFormat.DOCX
    assert parsed.detected_title == "ETCH-03 Maintenance"
    assert len(parsed.tables) == 1
    assert len(parsed.images) == 1
    assert payload_store.read(parsed.images[0].payload).startswith(b"\x89PNG")
    chunk_types = [chunk.chunk_type for chunk in parsed.chunks]
    assert ChunkType.TABLE in chunk_types
    assert ChunkType.IMAGE_TEXT in chunk_types
    assert any(chunk.title_path == ("Chamber Clean",) for chunk in parsed.chunks)
    assert parsed.images[0].related_chunk_draft_ids
    _assert_frozen_contract(parsed)


def test_xlsx_splits_regions_and_records_hidden_formula_merge_and_units() -> None:
    parsed = build_dispatcher(ProcessPayloadStore()).parse(
        "fdc.xlsx",
        _xlsx_bytes(),
        declared_media_type=("application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"),
        correlation_id="xlsx-1",
    )

    assert parsed.source_format is SourceFormat.XLSX
    assert parsed.metrics.sheets == 1
    assert len(parsed.tables) == 2
    warning_codes = {warning.code for warning in parsed.warnings}
    assert "INGEST_WARNING_HIDDEN_SHEETS_SKIPPED" in warning_codes
    assert "INGEST_WARNING_FORMULA_CACHE_MISSING" in warning_codes
    assert parsed.chunks[0].metadata["units"] == {
        "Pressure(kPa)": "kPa",
        "RF Match(%)": "%",
    }
    assert any("A5:B5" in chunk.metadata["merged_ranges"] for chunk in parsed.chunks)
    assert all(chunk.location.sheet_name == "FDC" for chunk in parsed.chunks)
    _assert_frozen_contract(parsed)


def test_pptx_preserves_slide_text_table_image_and_bbox() -> None:
    payload_store = ProcessPayloadStore()
    parsed = build_dispatcher(payload_store).parse(
        "inspection.pptx",
        _pptx_bytes(),
        declared_media_type=(
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ),
        correlation_id="pptx-1",
    )

    assert parsed.source_format is SourceFormat.PPTX
    assert parsed.metrics.slides == 1
    assert parsed.detected_title == "Wafer Edge Defect"
    assert len(parsed.tables) == 1
    assert len(parsed.images) == 1
    assert any(chunk.location.slide_number == 1 for chunk in parsed.chunks)
    assert all(
        chunk.location.bbox is not None
        for chunk in parsed.chunks
        if chunk.chunk_type in {ChunkType.TABLE, ChunkType.IMAGE_TEXT}
    )
    _assert_frozen_contract(parsed)
